import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def main():
    prs = Presentation()
    # Set 16:9 widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Scheme (Premium Dark Tech Mode)
    BG_COLOR = RGBColor(10, 17, 40)       # Dark Navy
    CARD_COLOR = RGBColor(20, 27, 54)     # Slightly lighter dark navy for content cards
    TITLE_COLOR = RGBColor(56, 189, 248)  # Bright Sky Blue
    TEXT_COLOR = RGBColor(241, 245, 249)  # Off-white
    MUTED_COLOR = RGBColor(148, 163, 184) # Muted Silver-Gray
    ACCENT_COLOR = RGBColor(45, 212, 191) # Mint/Teal
    CODE_BG = RGBColor(15, 23, 42)        # Very dark gray/black for code blocks
    CARD_BORDER = RGBColor(30, 41, 73)

    blank_layout = prs.slide_layouts[6]

    # Helper function to paint slide background
    def paint_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    # Helper function to create content cards
    def add_card(slide, left, top, width, height, title_text=None, title_color=ACCENT_COLOR):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        if title_text:
            tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Segoe UI'
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = title_color
            
        return card

    # Helper to add standard header
    def add_slide_header(slide, title_text, category="BIZIONARY ERP SYSTEM"):
        paint_bg(slide)
        
        # Category label
        cat_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.3))
        cat_tf = cat_tb.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_right = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = 'Segoe UI'
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = ACCENT_COLOR
        
        # Main Title
        title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(12.333), Inches(0.8))
        title_tf = title_tb.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = 'Segoe UI'
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = TITLE_COLOR

    # Helper to add bullet points (Easy English, brief)
    def add_bullet_points(slide, items, left, top, width, height, font_size=13):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        for idx, item in enumerate(items):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.font.name = 'Segoe UI'
            p.font.size = Pt(font_size)
            p.space_after = Pt(4)
            
            if ': ' in item and not item.startswith('http'):
                parts = item.split(': ', 1)
                run1 = p.add_run()
                run1.text = "• " + parts[0] + ": "
                run1.font.bold = True
                run1.font.color.rgb = TITLE_COLOR
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.bold = False
                run2.font.color.rgb = TEXT_COLOR
            else:
                p.text = "• " + item
                p.font.color.rgb = TEXT_COLOR

    # Helper to insert UI image with card outline
    def add_ui_image(slide, image_filename, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.04), top - Inches(0.04), width + Inches(0.08), height + Inches(0.08))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.color.rgb = ACCENT_COLOR
        card.line.width = Pt(1.5)
        
        image_path = os.path.join(os.getcwd(), "ui pics", image_filename)
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            tb = slide.shapes.add_textbox(left, top + (height/2) - Inches(0.5), width, Inches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[Image Missing: {image_filename}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Segoe UI'
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(239, 68, 68)

    # Helper to draw horizontal flow diagram at the bottom of Right Card
    def add_flow_diagram(slide, left, top, steps):
        x = left
        for i, step in enumerate(steps):
            # Draw Step Rounded Rect
            block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, Inches(1.05), Inches(0.55))
            block.fill.solid()
            block.fill.fore_color.rgb = CODE_BG
            block.line.color.rgb = ACCENT_COLOR
            block.line.width = Pt(1.0)
            
            # Text block inside
            tf = block.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            p = tf.paragraphs[0]
            p.text = step
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Segoe UI'
            p.font.size = Pt(8.5)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR
            
            x += Inches(1.05)
            
            # Draw arrow if not last step
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(0.05), top + Inches(0.18), Inches(0.18), Inches(0.2))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = TITLE_COLOR
                arrow.line.fill.background()
                x += Inches(0.28)

    # ==========================================
    # SLIDE 1: Title & Introduction Slide
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    paint_bg(s1)
    
    tb = s1.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BIZIONARY ERP SYSTEM"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = "A Secure, Agentic AI-Enabled Enterprise Resource Planning Platform for SMEs"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_COLOR
    p2.space_before = Pt(8)
    
    # Bottom columns
    add_card(s1, Inches(0.75), Inches(4.8), Inches(3.7), Inches(1.8), "Corporate Ledgers")
    points_1 = [
        "Relational database schemas.",
        "Double-entry ledger logs.",
        "Automatic stock signals."
    ]
    add_bullet_points(s1, points_1, Inches(0.95), Inches(5.3), Inches(3.3), Inches(1.2), font_size=12)
    
    add_card(s1, Inches(4.8), Inches(4.8), Inches(3.7), Inches(1.8), "AI-Driven Insights")
    points_2 = [
        "Groq Llama 3.3 chatbot.",
        "NLP pricing recommendations.",
        "Stock velocity demand forecast."
    ]
    add_bullet_points(s1, points_2, Inches(5.0), Inches(5.3), Inches(3.3), Inches(1.2), font_size=12)
    
    add_card(s1, Inches(8.85), Inches(4.8), Inches(3.7), Inches(1.8), "Dynamic Ingest")
    points_3 = [
        "Drag & Drop sales PDF upload.",
        "Monthly Excel parser.",
        "Zero-downtime key rotation."
    ]
    add_bullet_points(s1, points_3, Inches(9.05), Inches(5.3), Inches(3.3), Inches(1.2), font_size=12)

    # ==========================================
    # SLIDE 2: Problem & Solution
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_header(s2, "The Problem vs. The Bizionary Solution")
    
    add_card(s2, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8), "SME Operational Bottlenecks", RGBColor(239, 68, 68))
    prob_points = [
        "Data Silos: Using paper bills and offline sheets causes inventory mismatch.",
        "No Auditing: Database changes are done without logs, leading to errors.",
        "No Real-time Profit View: Creating financial statements takes weeks.",
        "Bad Restock Habits: Checking stock by hand leads to out-of-stock items."
    ]
    add_bullet_points(s2, prob_points, Inches(0.95), Inches(2.4), Inches(5.2), Inches(3.9), font_size=14)
    
    add_card(s2, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), "The Integrated Solution", ACCENT_COLOR)
    sol_points = [
        "One Database: Connects Sales, Purchases, Ledgers, and stock in one place.",
        "Auto Journals: Database signals save debit/credit journals automatically.",
        "Sub-Second AI Help: Chatbot queries the live DB to show reports.",
        "Auto-Ingest: Drag-and-drop Excel or PDFs to update inventory quickly."
    ]
    add_bullet_points(s2, sol_points, Inches(7.18), Inches(2.4), Inches(5.2), Inches(3.9), font_size=14)

    # ==========================================
    # SLIDE 3: Tech Stack (Visualization Diagram)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_header(s3, "Enterprise Technology Stack")
    
    # Grid of card blocks
    add_card(s3, Inches(0.5), Inches(1.8), Inches(2.2), Inches(4.8), "React Client", TITLE_COLOR)
    fe_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.4), Inches(1.9), Inches(3.9))
    fe_rect.fill.solid(); fe_rect.fill.fore_color.rgb = CODE_BG; fe_rect.line.color.rgb = CARD_BORDER
    tf1 = fe_rect.text_frame; tf1.word_wrap = True
    p1 = tf1.paragraphs[0]; p1.text = "FRONTEND SPA\n\n• React 19.2 (Vite)\n• Tailwind CSS v4\n• Lucide Icons\n• Recharts graphs\n• jsPDF downloads\n• HTML Canvas"
    p1.font.name = 'Segoe UI'; p1.font.size = Pt(12); p1.font.color.rgb = TEXT_COLOR
    
    add_card(s3, Inches(3.0), Inches(1.8), Inches(2.2), Inches(4.8), "Django Server", ACCENT_COLOR)
    be_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.15), Inches(2.4), Inches(1.9), Inches(3.9))
    be_rect.fill.solid(); be_rect.fill.fore_color.rgb = CODE_BG; be_rect.line.color.rgb = CARD_BORDER
    tf2 = be_rect.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = "APPLICATION CORE\n\n• Django 4.2.7 Core\n• REST Framework\n• JWT Stateless Auth\n• post_save Signals\n• CORS Middleware"
    p2.font.name = 'Segoe UI'; p2.font.size = Pt(12); p2.font.color.rgb = TEXT_COLOR
    
    add_card(s3, Inches(5.5), Inches(1.8), Inches(2.2), Inches(4.8), "Cognitive AI", RGBColor(168, 85, 247))
    ai_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.65), Inches(2.4), Inches(1.9), Inches(3.9))
    ai_rect.fill.solid(); ai_rect.fill.fore_color.rgb = CODE_BG; ai_rect.line.color.rgb = CARD_BORDER
    tf3 = ai_rect.text_frame; tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.text = "INTELLIGENT LAYER\n\n• Groq SDK\n• Llama 3.3 Engine\n• Function Calling\n• OpenAI API GPT\n• Sentiment Analysis\n• Predictive Insights"
    p3.font.name = 'Segoe UI'; p3.font.size = Pt(12); p3.font.color.rgb = TEXT_COLOR
    
    add_card(s3, Inches(8.0), Inches(1.8), Inches(2.2), Inches(4.8), "Data & Ingestion", RGBColor(234, 179, 8))
    dt_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.15), Inches(2.4), Inches(1.9), Inches(3.9))
    dt_rect.fill.solid(); dt_rect.fill.fore_color.rgb = CODE_BG; dt_rect.line.color.rgb = CARD_BORDER
    tf4 = dt_rect.text_frame; tf4.word_wrap = True
    p4 = tf4.paragraphs[0]; p4.text = "ANALYSIS & STORAGE\n\n• SQLite 3 (Local)\n• PostgreSQL (Cloud)\n• Pandas parser\n• openpyxl Engine\n• DB Cache manager"
    p4.font.name = 'Segoe UI'; p4.font.size = Pt(12); p4.font.color.rgb = TEXT_COLOR
    
    add_card(s3, Inches(10.5), Inches(1.8), Inches(2.33), Inches(4.8), "Infrastructure", RGBColor(239, 68, 68))
    inf_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.65), Inches(2.4), Inches(2.03), Inches(3.9))
    inf_rect.fill.solid(); inf_rect.fill.fore_color.rgb = CODE_BG; inf_rect.line.color.rgb = CARD_BORDER
    tf5 = inf_rect.text_frame; tf5.word_wrap = True
    p5 = tf5.paragraphs[0]; p5.text = "DEPLOYMENT\n\n• Vercel Edge CDN\n• Railway API host\n• Docker runtimes\n• PyInstaller build\n• Windows standalone"
    p5.font.name = 'Segoe UI'; p5.font.size = Pt(12); p5.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 4: Whole Project System Architecture (Detailed Visualization Diagram)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_header(s4, "Project Components & Data Flow Diagram")
    
    # helper for drawing small items
    def add_sub_item(slide, left, top, width, height, text, bg_color):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_color
        rect.line.color.rgb = CARD_BORDER
        rect.line.width = Pt(1.0)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR
        return rect

    # 1. Frontend Client (React SPA)
    add_card(s4, Inches(0.5), Inches(2.0), Inches(2.6), Inches(4.8), "Client (React SPA)", TITLE_COLOR)
    add_sub_item(s4, Inches(0.7), Inches(2.7), Inches(2.2), Inches(0.7), "Vite Single Page UI\n(Responsive Screens)", CODE_BG)
    add_sub_item(s4, Inches(0.7), Inches(3.6), Inches(2.2), Inches(0.7), "Custom State Hooks\n(Axios Requests)", CODE_BG)
    add_sub_item(s4, Inches(0.7), Inches(4.5), Inches(2.2), Inches(0.7), "Recharts Visualization\n(Sales Charts)", CODE_BG)
    add_sub_item(s4, Inches(0.7), Inches(5.4), Inches(2.2), Inches(0.7), "jsPDF Invoicing\n(Offline PDF exports)", CODE_BG)
    
    # Arrow 1: FE -> BE (Right Arrow)
    a1 = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(4.1), Inches(1.1), Inches(0.4))
    a1.fill.solid(); a1.fill.fore_color.rgb = ACCENT_COLOR; a1.line.fill.background()

    # 2. Backend Server (Django REST)
    add_card(s4, Inches(4.4), Inches(2.0), Inches(2.6), Inches(4.8), "Server (Django REST)", ACCENT_COLOR)
    add_sub_item(s4, Inches(4.6), Inches(2.7), Inches(2.2), Inches(0.7), "API Router Gateway\n(JWT Authentication)", CODE_BG)
    add_sub_item(s4, Inches(4.6), Inches(3.6), Inches(2.2), Inches(0.7), "post_save Signals\n(Ledger post hooks)", CODE_BG)
    add_sub_item(s4, Inches(4.6), Inches(4.5), Inches(2.2), Inches(0.7), "Excel Pandas Parser\n(Bulk synchronizer)", CODE_BG)
    add_sub_item(s4, Inches(4.6), Inches(5.4), Inches(2.2), Inches(0.7), "API Key Cache Manager\n(In-memory caching)", CODE_BG)

    # Connecting BE to DB and AI Tiers
    # Arrow 2: BE -> DB (Right Arrow top)
    a2 = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.1), Inches(3.0), Inches(1.1), Inches(0.3))
    a2.fill.solid(); a2.fill.fore_color.rgb = TITLE_COLOR; a2.line.fill.background()
    
    # Arrow 3: DB -> BE (Left Arrow top)
    a3 = s4.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.1), Inches(3.6), Inches(1.1), Inches(0.3))
    a3.fill.solid(); a3.fill.fore_color.rgb = MUTED_COLOR; a3.line.fill.background()

    # Arrow 4: BE -> AI (Right Arrow bottom)
    a4 = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.1), Inches(4.8), Inches(1.1), Inches(0.3))
    a4.fill.solid(); a4.fill.fore_color.rgb = TITLE_COLOR; a4.line.fill.background()
    
    # Arrow 5: AI -> BE (Left Arrow bottom)
    a5 = s4.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(7.1), Inches(5.4), Inches(1.1), Inches(0.3))
    a5.fill.solid(); a5.fill.fore_color.rgb = MUTED_COLOR; a5.line.fill.background()

    # 3. Database Tier (Relational DB)
    add_card(s4, Inches(8.3), Inches(2.0), Inches(2.2), Inches(4.8), "Database Tier", RGBColor(234, 179, 8))
    add_sub_item(s4, Inches(8.5), Inches(2.8), Inches(1.8), Inches(1.2), "SQLite / Postgres\n\n(Enforces integrity,\nforeign key constraints)", CODE_BG)
    add_sub_item(s4, Inches(8.5), Inches(4.6), Inches(1.8), Inches(1.2), "Double-Entry Journals\n\n(Balanced assets,\ncash flows, COGS)", CODE_BG)

    # 4. Cognitive Tier (External AI)
    add_card(s4, Inches(10.6), Inches(2.0), Inches(2.2), Inches(4.8), "Cognitive AI Tier", RGBColor(168, 85, 247))
    add_sub_item(s4, Inches(10.8), Inches(2.8), Inches(1.8), Inches(1.2), "Groq Llama 3.3\n\n(High-speed RAG,\nfunction calling)", CODE_BG)
    add_sub_item(s4, Inches(10.8), Inches(4.6), Inches(1.8), Inches(1.2), "OpenAI API\n\n(PDF text parsing,\nsentiment metrics)", CODE_BG)

    # ==========================================
    # SLIDE 4A: Backend Architecture & Database Schema
    # ==========================================
    s4a = prs.slides.add_slide(blank_layout)
    add_slide_header(s4a, "Data Links & Storage", "BIZIONARY BACKEND SERVICE")
    
    # 30% Theory Left Card
    add_card(s4a, Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.9), "How Data is Linked", TITLE_COLOR)
    s4a_points = [
        "Automated Links: Connects products, sales, and purchases so we don't have to enter them twice.",
        "Safe Deletes: Prevents accidental deletions from breaking our inventory history.",
        "Easy Storage: Works with simple local files or scales to cloud databases."
    ]
    add_bullet_points(s4a, s4a_points, Inches(0.7), Inches(2.4), Inches(4.1), Inches(4.0), font_size=13)
    
    # 70% Visualization Right Card (Database ERD Tables layout)
    add_card(s4a, Inches(5.5), Inches(1.8), Inches(7.33), Inches(4.9), "How Our Data is Connected", ACCENT_COLOR)
    
    # Draw Table Cards
    # Table 1: Product
    add_sub_item(s4a, Inches(5.8), Inches(2.3), Inches(3.0), Inches(1.3), "Products Table (Main)\n------------------\n• product_id\n• product_name\n• current_stock\n• shop_stock", CODE_BG)
    # Table 2: InventoryTransaction (FK to Product)
    add_sub_item(s4a, Inches(9.4), Inches(2.3), Inches(3.0), Inches(1.3), "Stock Logs Table\n------------------\n• log_id\n• product_id\n• change_amount\n• action_type (IN/OUT)", CODE_BG)
    # Table 3: Sale (FK to Product)
    add_sub_item(s4a, Inches(5.8), Inches(4.4), Inches(3.0), Inches(1.8), "Sales Table\n------------------\n• sale_id\n• product_id\n• quantity_sold\n• total_cost\n• paid_status", CODE_BG)
    # Table 4: CashTransaction
    add_sub_item(s4a, Inches(9.4), Inches(4.4), Inches(3.0), Inches(1.8), "Cash Book Table\n------------------\n• transaction_id\n• link_id\n• link_type\n• amount\n• type (IN/OUT)", CODE_BG)

    # Connections
    # Product Table -> InventoryTransaction Table
    a_prod_inv = s4a.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.95), Inches(2.8), Inches(0.35), Inches(0.25))
    a_prod_inv.fill.solid(); a_prod_inv.fill.fore_color.rgb = TITLE_COLOR; a_prod_inv.line.fill.background()
    
    # Sale Table -> CashTransaction Table
    a_sale_cash = s4a.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.95), Inches(5.15), Inches(0.35), Inches(0.25))
    a_sale_cash.fill.solid(); a_sale_cash.fill.fore_color.rgb = ACCENT_COLOR; a_sale_cash.line.fill.background()

    # ==========================================
    # SLIDE 4B: API Communication & Security Pipeline
    # ==========================================
    s4b = prs.slides.add_slide(blank_layout)
    add_slide_header(s4b, "Safe Screen-to-Server Talk", "BIZIONARY BACKEND SERVICE")
    
    # 30% Theory Left Card
    add_card(s4b, Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.9), "Safe Data Travel", TITLE_COLOR)
    s4b_points = [
        "Secure Login Keys: Users get a unique login token to check invoices securely.",
        "Security Gates: Blocks hackers and blocks unsafe websites from accessing data.",
        "Data Validation: Scans every input form instantly and blocks bad values."
    ]
    add_bullet_points(s4b, s4b_points, Inches(0.7), Inches(2.4), Inches(4.1), Inches(4.0), font_size=13)
    
    # 70% Visualization Right Card (API Loop Layout)
    add_card(s4b, Inches(5.5), Inches(1.8), Inches(7.33), Inches(4.9), "API Call Lifecycle (Clockwise Flow)", ACCENT_COLOR)
    
    # Clockwise API cycle blocks
    add_sub_item(s4b, Inches(5.8), Inches(2.3), Inches(3.0), Inches(0.9), "1. Web Screen\n(User clicks 'Save Sale'\non their browser)", CODE_BG)
    add_sub_item(s4b, Inches(9.4), Inches(2.3), Inches(3.0), Inches(0.9), "2. Security Check\n(Checks if user has a\nvalid security login key)", CODE_BG)
    add_sub_item(s4b, Inches(9.4), Inches(3.7), Inches(3.0), Inches(0.9), "3. Input Scanner\n(Checks if fields are filled\ncorrectly without errors)", CODE_BG)
    add_sub_item(s4b, Inches(5.8), Inches(3.7), Inches(3.0), Inches(0.9), "4. Business Logic\n(Deducts stock quantities\nand calculates cash)", CODE_BG)
    add_sub_item(s4b, Inches(5.8), Inches(5.1), Inches(6.6), Inches(0.9), "5. Database File\n(Saves transaction rows atomically and sends back a success message)", CODE_BG)
    
    # Connections for clockwise flow
    # React Client -> JWT Gate (Right)
    a_f1 = s4b.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.9), Inches(2.6), Inches(0.4), Inches(0.3))
    a_f1.fill.solid(); a_f1.fill.fore_color.rgb = TITLE_COLOR; a_f1.line.fill.background()
    # JWT Gate -> DRF Serializer (Down)
    a_f2 = s4b.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.75), Inches(3.3), Inches(0.3), Inches(0.3))
    a_f2.fill.solid(); a_f2.fill.fore_color.rgb = TITLE_COLOR; a_f2.line.fill.background()
    # DRF Serializer -> Django View (Left)
    a_f3 = s4b.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.9), Inches(4.0), Inches(0.4), Inches(0.3))
    a_f3.fill.solid(); a_f3.fill.fore_color.rgb = ACCENT_COLOR; a_f3.line.fill.background()
    # Django View -> DB Layer (Down)
    a_f4 = s4b.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.15), Inches(4.7), Inches(0.3), Inches(0.3))
    a_f4.fill.solid(); a_f4.fill.fore_color.rgb = ACCENT_COLOR; a_f4.line.fill.background()

    # ==========================================
    # SLIDE 4C: Signal-Driven Ledger & Auditing Pipeline
    # ==========================================
    s4c = prs.slides.add_slide(blank_layout)
    add_slide_header(s4c, "Automatic Bookkeeping Flow", "BIZIONARY BACKEND SERVICE")
    
    # 30% Theory Left Card
    add_card(s4c, Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.9), "Automated Ledgers", TITLE_COLOR)
    s4c_points = [
        "One-Click Updates: Saving a sale updates stock, logs cash, and writes audits automatically.",
        "All-or-Nothing: If one log fails, the entire sale cancels to keep numbers correct.",
        "Auto Reverse: Deleting a sale automatically returns stock and refunds cash."
    ]
    add_bullet_points(s4c, s4c_points, Inches(0.7), Inches(2.4), Inches(4.1), Inches(4.0), font_size=13)
    
    # 70% Visualization Right Card (Signal diagram)
    add_card(s4c, Inches(5.5), Inches(1.8), Inches(7.33), Inches(4.9), "Auto-Ledger & Bookkeeping Flow", ACCENT_COLOR)
    
    # Flow elements
    # Root
    add_sub_item(s4c, Inches(7.66), Inches(2.2), Inches(3.0), Inches(0.8), "New Transaction Saved\n(e.g., Sale is logged)", CODE_BG)
    
    # 3 Parallel Branches
    add_sub_item(s4c, Inches(5.8), Inches(3.8), Inches(2.0), Inches(1.1), "Deduct Inventory\n(Removes item\nfrom stock levels)", CODE_BG)
    add_sub_item(s4c, Inches(8.15), Inches(3.8), Inches(2.0), Inches(1.1), "Log Cash Inflow\n(Records cash\nreceived)", CODE_BG)
    add_sub_item(s4c, Inches(10.5), Inches(3.8), Inches(2.0), Inches(1.1), "Write Activity Log\n(Saves who did what\nfor auditing)", CODE_BG)
    
    # Atomic barrier
    add_sub_item(s4c, Inches(5.8), Inches(5.4), Inches(6.7), Inches(0.8), "Atomic Transaction Guard\n[All steps must succeed together, or everything is undone]", CODE_BG)

    # Branch down-arrows from root to child blocks
    a_b1 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.7), Inches(3.15), Inches(0.2), Inches(0.55))
    a_b1.fill.solid(); a_b1.fill.fore_color.rgb = TITLE_COLOR; a_b1.line.fill.background()
    
    a_b2 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.15), Inches(3.15), Inches(0.2), Inches(0.55))
    a_b2.fill.solid(); a_b2.fill.fore_color.rgb = TITLE_COLOR; a_b2.line.fill.background()
    
    a_b3 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.5), Inches(3.15), Inches(0.2), Inches(0.55))
    a_b3.fill.solid(); a_b3.fill.fore_color.rgb = TITLE_COLOR; a_b3.line.fill.background()

    # Branch down-arrows from child blocks to barrier
    a_c1 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.7), Inches(5.05), Inches(0.18), Inches(0.28))
    a_c1.fill.solid(); a_c1.fill.fore_color.rgb = ACCENT_COLOR; a_c1.line.fill.background()
    
    a_c2 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.15), Inches(5.05), Inches(0.18), Inches(0.28))
    a_c2.fill.solid(); a_c2.fill.fore_color.rgb = ACCENT_COLOR; a_c2.line.fill.background()
    
    a_c3 = s4c.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.5), Inches(5.05), Inches(0.18), Inches(0.28))
    a_c3.fill.solid(); a_c3.fill.fore_color.rgb = ACCENT_COLOR; a_c3.line.fill.background()

    # ==========================================
    # SLIDE 5: Raw Master Catalog (Al-Noor Trading)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_header(s5, "Raw Product Master Catalog Data")
    add_ui_image(s5, "media__1784459864365.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s5, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Product Catalog Excel Sheet", TITLE_COLOR)
    s5_data_points = [
        "Store Master Data: Defines Lahore office product keys, categorizations, and brands.",
        "Key Fields: Includes SKUs, reorder limits, cost values, selling margins, and supplier listings.",
        "System Mapping: The raw Excel contains structural attributes that are parsed into SQL product objects."
    ]
    add_bullet_points(s5, s5_data_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(2.0), font_size=13)
    add_flow_diagram(s5, Inches(8.85), Inches(4.5), ["1. Open Excel", "2. Map Schema", "3. Seed DB"])

    # ==========================================
    # SLIDE 6: Raw 30-Day Sales Data (Al-Noor Trading)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_header(s6, "Raw 30-Day Monthly Sales Data Sheet")
    add_ui_image(s6, "media__1784459858013.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s6, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Sales Sheet Columns", ACCENT_COLOR)
    s6_data_points = [
        "Monthly Logs: Stores transaction rows containing margins, reorder units, and stock statuses.",
        "Daily Columns: Tracks unit sales quantities sold day-by-day (from 1-May to 30-May).",
        "Parser Connection: Scans date headers dynamically to index quantities without hardcoded months."
    ]
    add_bullet_points(s6, s6_data_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(2.0), font_size=13)
    add_flow_diagram(s6, Inches(8.85), Inches(4.5), ["1. Read Columns", "2. Parse Quantities", "3. Record Sales"])

    # ==========================================
    # SLIDE 7: Master Data & 6-Month Sales Seeding Flow
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_header(s7, "Master Data & 6-Month Sales Seeding")
    
    # Sequence boxes (3 steps)
    add_card(s7, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), "Step 1: Master Catalog Data", TITLE_COLOR)
    st1 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.5), Inches(3.5), Inches(3.8))
    st1.fill.solid(); st1.fill.fore_color.rgb = CODE_BG; st1.line.color.rgb = CARD_BORDER
    t_st1 = st1.text_frame; t_st1.word_wrap = True
    p_st1 = t_st1.paragraphs[0]; p_st1.text = "Initial Data Setup:\n• Exposes real catalogs of items, standard purchase costs, supplier profiles, and categories.\n• Registers items like Beverages, Grocery, Pharmaceuticals, and Stationery.\n• Pre-defines minimum stock thresholds for alert triggers."
    p_st1.font.name = 'Segoe UI'; p_st1.font.size = Pt(13); p_st1.font.color.rgb = TEXT_COLOR
    
    arrow_st1 = s7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.45), Inches(3.8), Inches(0.4), Inches(0.3))
    arrow_st1.fill.solid(); arrow_st1.fill.fore_color.rgb = ACCENT_COLOR; arrow_st1.line.fill.background()

    add_card(s7, Inches(5.0), Inches(2.0), Inches(3.8), Inches(4.5), "Step 2: Seeding Engine", ACCENT_COLOR)
    st2 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.15), Inches(2.5), Inches(3.5), Inches(3.8))
    st2.fill.solid(); st2.fill.fore_color.rgb = CODE_BG; st2.line.color.rgb = CARD_BORDER
    t_st2 = st2.text_frame; t_st2.word_wrap = True
    p_st2 = t_st2.paragraphs[0]; p_st2.text = "Seeding Logic:\n• Executes script `seed_historical_sales.py` to generate sample logs.\n• Seeds data across multiple calendar months (e.g. January to April 2026).\n• Simulates random order sizes (1-8 pcs) and standard payment methods."
    p_st2.font.name = 'Segoe UI'; p_st2.font.size = Pt(13); p_st2.font.color.rgb = TEXT_COLOR
    
    arrow_st2 = s7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.95), Inches(3.8), Inches(0.4), Inches(0.3))
    arrow_st2.fill.solid(); arrow_st2.fill.fore_color.rgb = ACCENT_COLOR; arrow_st2.line.fill.background()

    add_card(s7, Inches(9.5), Inches(2.0), Inches(3.33), Inches(4.5), "Step 3: Outcome & Analytics", RGBColor(234, 179, 8))
    st3 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.65), Inches(2.5), Inches(3.03), Inches(3.8))
    st3.fill.solid(); st3.fill.fore_color.rgb = CODE_BG; st3.line.color.rgb = CARD_BORDER
    t_st3 = st3.text_frame; t_st3.word_wrap = True
    p_st3 = t_st3.paragraphs[0]; p_st3.text = "Outcome Results:\n• Generates 6 months of sales history representing over 100+ transactions.\n• Populates dashboard insights with real trending sales velocity curves.\n• Enables AI analytics to test demand forecasts, profit summaries, and reorder levels."
    p_st3.font.name = 'Segoe UI'; p_st3.font.size = Pt(13); p_st3.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 8: Executive Dashboard (Screenshot 012114.png)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_header(s8, "Executive Dashboard Overview")
    add_ui_image(s8, "Screenshot 2026-07-19 012114.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s8, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Executive Summary stats", TITLE_COLOR)
    s8_points = [
        "Overview: Displays total sales, expenses, net profits, and inventory value.",
        "Quick Buttons: Easily add products, register suppliers, or log adjustments."
    ]
    add_bullet_points(s8, s8_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    # Graphic Flow Diagram
    add_flow_diagram(s8, Inches(8.85), Inches(4.5), ["1. Open Screen", "2. Fetch Stats", "3. Show Dashboard"])

    # ==========================================
    # SLIDE 9: Sales Insights Dashboard (Screenshot 012111.png)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_header(s9, "Sales Performance Insights Dashboard")
    add_ui_image(s9, "Screenshot 2026-07-19 012111.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s9, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Interactive Sales Charts", ACCENT_COLOR)
    s9_points = [
        "Charts: Shows categories sold (stacked bars) and revenue trends (lines).",
        "Filtering: Toggle days to filter trends instantly."
    ]
    add_bullet_points(s9, s9_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s9, Inches(8.85), Inches(4.5), ["1. Pick Period", "2. Group Dates", "3. Recharts Graph"])

    # ==========================================
    # SLIDE 10: Accounts & Finance Ledger (Screenshot 012108.png)
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_header(s10, "Accounts & Financial Ledger Module")
    add_ui_image(s10, "Screenshot 2026-07-19 012108.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s10, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Double-Entry Audits", TITLE_COLOR)
    s10_points = [
        "Financial Tabs: View Revenues, Expenses, Utility Bills, and statement logs.",
        "Reconciliation: Simple button to verify debits equal credits instantly."
    ]
    add_bullet_points(s10, s10_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s10, Inches(8.85), Inches(4.5), ["1. Log Action", "2. Signals Trigger", "3. Ledger Entry"])

    # ==========================================
    # SLIDE 11: Product Catalog Grid (Screenshot 012102.png)
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_header(s11, "Product Catalog & Custom Sections")
    add_ui_image(s11, "Screenshot 2026-07-19 012102.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s11, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Inventory Catalog System", ACCENT_COLOR)
    s11_points = [
        "Grid List: Shows catalog product SKUs, prices, margins, and active status.",
        "Custom Columns: Create extra fields for specific categories dynamically."
    ]
    add_bullet_points(s11, s11_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s11, Inches(8.85), Inches(4.5), ["1. Add Columns", "2. Alter Schema", "3. Update Tables"])

    # ==========================================
    # SLIDE 12: Stock Management Dashboard (Screenshot 012055.png)
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_header(s12, "Stock Management Control Center")
    add_ui_image(s12, "Screenshot 2026-07-19 012055.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s12, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Reorder Alerts System", TITLE_COLOR)
    s12_points = [
        "Valuations: Tracks shop retail value vs. warehouse stock valuations.",
        "Alerts: Highlights items where stock quantity is below safety limits."
    ]
    add_bullet_points(s12, s12_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s12, Inches(8.85), Inches(4.5), ["1. Check Stock", "2. Compare Limit", "3. Trigger Alert"])

    # ==========================================
    # SLIDE 13: Warehouse & Incoming Stock Modals (Screenshot 012058.png & 012050.png)
    # ==========================================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_header(s13, "Warehouse Stock & Incoming Modals")
    add_ui_image(s13, "Screenshot 2026-07-19 012058.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s13, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Stock Locations Reports", ACCENT_COLOR)
    s13_points = [
        "Location Splitting: Traces items inside the warehouse vs. items in shop.",
        "Incoming Orders: List of pending items ordered from supplier shipments."
    ]
    add_bullet_points(s13, s13_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s13, Inches(8.85), Inches(4.5), ["1. Select Item", "2. Query Location", "3. Display Modal"])

    # ==========================================
    # SLIDE 14: Sales Transaction Log & Ingestion Suite (Screenshot 012046.png)
    # ==========================================
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_header(s14, "Sales Transactions & Ingestion Suite")
    add_ui_image(s14, "Screenshot 2026-07-19 012046.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s14, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Transactions Logging", TITLE_COLOR)
    s14_points = [
        "Logs: Search and filter complete logs of customer sales and margins.",
        "Auto-Ingest: Drag-and-drop Excel sheets to parse transactions quickly."
    ]
    add_bullet_points(s14, s14_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s14, Inches(8.85), Inches(4.5), ["1. Import Sheet", "2. Pandas Parser", "3. Log Sale DB"])

    # ==========================================
    # SLIDE 15: Supplier Ordered Slips (Screenshot 012037.png)
    # ==========================================
    s15 = prs.slides.add_slide(blank_layout)
    add_slide_header(s15, "Supplier Ordered Slips Management")
    add_ui_image(s15, "Screenshot 2026-07-19 012037.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s15, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Procurement Cycles", ACCENT_COLOR)
    s15_points = [
        "Tracking: Monitor supplier slips, quantities, total sums, and pending statuses.",
        "Stock-in: Easily mark items as received to update inventories immediately."
    ]
    add_bullet_points(s15, s15_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s15, Inches(8.85), Inches(4.5), ["1. Mark Received", "2. Update Slip Ratio", "3. Inflow Stock"])

    # ==========================================
    # SLIDE 16: AI Chatbot Assistant Interface (Screenshot 012026.png)
    # ==========================================
    s16 = prs.slides.add_slide(blank_layout)
    add_slide_header(s16, "AI Chatbot Assistant Interface")
    add_ui_image(s16, "Screenshot 2026-07-19 012026.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s16, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Agentic AI Chatbot", TITLE_COLOR)
    s16_points = [
        "RAG Agent: Answers questions on invoices, stock alerts, and revenues.",
        "Quick Shortcuts: One-click buttons to check margins or display charts."
    ]
    add_bullet_points(s16, s16_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s16, Inches(8.85), Inches(4.5), ["1. User Prompt", "2. Groq Tool Use", "3. Local Python DB"])

    # ==========================================
    # SLIDE 17: Pack Price Calculator (media__1784456905899.png)
    # ==========================================
    s17 = prs.slides.add_slide(blank_layout)
    add_slide_header(s17, "Pack & Carton Price Calculator")
    add_ui_image(s17, "media__1784456905899.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s17, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Pack Dynamics Calculator", ACCENT_COLOR)
    s17_points = [
        "Inputs: Configure product single unit selling prices and carton multipliers.",
        "Profit Projections: Displays expected net revenue and margins instantly."
    ]
    add_bullet_points(s17, s17_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s17, Inches(8.85), Inches(4.5), ["1. Enter Quantity", "2. React State", "3. Compute Profit"])

    # ==========================================
    # SLIDE 18: Direct Stock Purchase (media__1784456944788.png)
    # ==========================================
    s18 = prs.slides.add_slide(blank_layout)
    add_slide_header(s18, "Direct Stock Purchase (Latest Modification)")
    add_ui_image(s18, "media__1784456944788.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s18, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "Direct Inflow Ledger Entries", TITLE_COLOR)
    s18_points = [
        "Stock Intake: Log supplier purchases directly to shop or warehouse stock.",
        "Custom Pack Sizes: Input unit cost and packs to update catalogs."
    ]
    add_bullet_points(s18, s18_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s18, Inches(8.85), Inches(4.5), ["1. Select Product", "2. Route Location", "3. Signal Stock In"])

    # ==========================================
    # SLIDE 19: AI Sales Slip PDF Upload (media__1784457199782.png)
    # ==========================================
    s19 = prs.slides.add_slide(blank_layout)
    add_slide_header(s19, "AI Sales Slip PDF Upload / Bulk Sale")
    add_ui_image(s19, "media__1784457199782.png", Inches(0.5), Inches(1.8), Inches(8.0), Inches(4.8))
    add_card(s19, Inches(8.7), Inches(1.8), Inches(4.13), Inches(4.8), "AI-Driven Parser Modal", ACCENT_COLOR)
    s19_points = [
        "Upload PDF: Extracts item name, quantities, and cost automatically.",
        "Fuzzy Matcher: AI maps scanned invoice details to existing catalog SKUs."
    ]
    add_bullet_points(s19, s19_points, Inches(8.85), Inches(2.4), Inches(3.8), Inches(1.5), font_size=13)
    
    add_flow_diagram(s19, Inches(8.85), Inches(4.5), ["1. Upload File", "2. OpenAI Extract", "3. Record Sales DB"])

    # ==========================================
    # SLIDE 20: Standalone Desktop Executable (.exe)
    # ==========================================
    s20 = prs.slides.add_slide(blank_layout)
    add_slide_header(s20, "Standalone Desktop Executable (.exe)")
    
    # 3 sequence blocks for compilation pipeline
    add_card(s20, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5), "Build Screen Files", TITLE_COLOR)
    exe_c1 = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.5), Inches(3.5), Inches(3.8))
    exe_c1.fill.solid(); exe_c1.fill.fore_color.rgb = CODE_BG; exe_c1.line.color.rgb = CARD_BORDER
    tf_c1 = exe_c1.text_frame; tf_c1.word_wrap = True
    p_c1 = tf_c1.paragraphs[0]; p_c1.text = "Step 1: Screen Code:\n• Packages interface code into clean, fast browser files.\n• Optimizes screens to load instantly on any user desktop.\n• Removes temporary folders to keep the app lightweight."
    p_c1.font.name = 'Segoe UI'; p_c1.font.size = Pt(13); p_c1.font.color.rgb = TEXT_COLOR
    
    arrow_c1 = s20.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.45), Inches(3.8), Inches(0.4), Inches(0.3))
    arrow_c1.fill.solid(); arrow_c1.fill.fore_color.rgb = ACCENT_COLOR; arrow_c1.line.fill.background()

    add_card(s20, Inches(5.0), Inches(2.0), Inches(3.8), Inches(4.5), "Prepare Database", ACCENT_COLOR)
    exe_c2 = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.15), Inches(2.5), Inches(3.5), Inches(3.8))
    exe_c2.fill.solid(); exe_c2.fill.fore_color.rgb = CODE_BG; exe_c2.line.color.rgb = CARD_BORDER
    tf_c2 = exe_c2.text_frame; tf_c2.word_wrap = True
    p_c2 = tf_c2.paragraphs[0]; p_c2.text = "Step 2: Database Setup:\n• Prepares the SQLite database file structures.\n• Sets up the tables and rules for products and invoices.\n• Collects final files using simple scripts."
    p_c2.font.name = 'Segoe UI'; p_c2.font.size = Pt(13); p_c2.font.color.rgb = TEXT_COLOR
    
    arrow_c2 = s20.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.95), Inches(3.8), Inches(0.4), Inches(0.3))
    arrow_c2.fill.solid(); arrow_c2.fill.fore_color.rgb = ACCENT_COLOR; arrow_c2.line.fill.background()

    add_card(s20, Inches(9.5), Inches(2.0), Inches(3.33), Inches(4.5), "Create Installer (.exe)", RGBColor(234, 179, 8))
    exe_c3 = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.65), Inches(2.5), Inches(3.03), Inches(3.8))
    exe_c3.fill.solid(); exe_c3.fill.fore_color.rgb = CODE_BG; exe_c3.line.color.rgb = CARD_BORDER
    tf_c3 = exe_c3.text_frame; tf_c3.word_wrap = True
    p_c3 = tf_c3.paragraphs[0]; p_c3.text = "Step 3: Standalone App:\n• Bundles server code and files into a double-clickable installer file.\n• Creates a single zip folder containing the offline runner.\n• Users can copy it to a USB drive and launch it on any Windows PC."
    p_c3.font.name = 'Segoe UI'; p_c3.font.size = Pt(13); p_c3.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 21: Data Portability & Disaster Recovery
    # ==========================================
    s21 = prs.slides.add_slide(blank_layout)
    add_slide_header(s21, "Data Backups & Upgrades")
    
    add_card(s21, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8), "All-in-One File Backups", TITLE_COLOR)
    port_points_1 = [
        "Easy Storage: All records (products, sales, cash logs) live in one file: `db.sqlite3`.",
        "Safe Location: The database file sits safely inside the app's internal folder.",
        "Copy-Paste Backup: Back up everything by simply making a copy of this single file."
    ]
    add_bullet_points(s21, port_points_1, Inches(1.0), Inches(2.4), Inches(5.1), Inches(3.8), font_size=16)
    
    add_card(s21, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), "How to Upgrade the App", ACCENT_COLOR)
    port_points_2 = [
        "1. Copy Old Database: Copy the `db.sqlite3` file from your current app folder.",
        "2. Unzip New Version: Unpack the new version folder you downloaded.",
        "3. Paste Database: Paste your copied file into the new app folder, replacing the blank one.",
        "4. Start App: Double-click the app icon to load all your records instantly."
    ]
    add_bullet_points(s21, port_points_2, Inches(7.23), Inches(2.4), Inches(5.1), Inches(3.8), font_size=16)

    # ==========================================
    # SLIDE 22: Future Enhancements
    # ==========================================
    s22 = prs.slides.add_slide(blank_layout)
    add_slide_header(s22, "Future Enhancements")
    
    add_card(s22, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.8), "Scalability & Product Roadmap", ACCENT_COLOR)
    future_points = [
        "Multi-Tenant SaaS Support: Transition the current single-enterprise design to support multi-tenant user bases, permitting multiple businesses to register and manage their operations independently.",
        "Hardware Integration: Build direct interfaces for POS peripheral devices, including barcode scanners, receipt thermal printers, and electronic cash drawers.",
        "Automated Notifications Pipeline: Connect SMS (Twilio) and Email (SendGrid) triggers to dynamically alert suppliers on low stock, send invoice reminders to customers, and email managers financial summaries.",
        "Extended AI Capabilities: Build deep learning models for advanced predictive demand forecasting based on external seasonal events, market trends, and historic data."
    ]
    add_bullet_points(s22, future_points, Inches(1.0), Inches(2.4), Inches(11.33), Inches(3.9), font_size=16)

    # ==========================================
    # SLIDE 23: Thank You & Q&A
    # ==========================================
    s23 = prs.slides.add_slide(blank_layout)
    paint_bg(s23)
    
    tb = s23.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.833), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Questions & Answers Session"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    
    # Save the presentation
    filename = "BizionaryERP_Presentation_v6.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == '__main__':
    main()
